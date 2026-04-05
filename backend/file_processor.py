"""
file_processor.py — Extract text from uploaded files.

Supported:
  .pdf          → pdfplumber
  .docx / .doc  → python-docx
  .pptx / .ppt  → python-pptx
  .txt / .md    → plain UTF-8
  .jpg / .jpeg / .png / .gif / .webp → base64 for vision LLM
  .mp4 / .mov / .avi / etc.          → not supported (returns guidance message)
"""

import base64
import io
import json
from pathlib import Path

# Reliability scores per user-supplied label
RELIABILITY_SCORES = {
    "official":  10,   # assignment brief, rubric, professor instructions
    "reference":  7,   # supplementary reading, example code
    "notes":      4,   # informal notes, rough drafts
    "auto":       6,   # unclassified (default)
}

IMAGE_EXTS  = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
TEXT_EXTS   = {".txt", ".md", ".csv"}
VIDEO_EXTS  = {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v"}


def extract_text(file_bytes: bytes, filename: str) -> tuple[str, str]:
    """
    Returns (content, file_type).
    For images: content is a JSON string {"type":"image","mime":...,"b64":...}
    For video:  content is a guidance message string, file_type = "video"
    For text:   content is plain extracted text, file_type = pdf|docx|pptx|text|image
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _pdf(file_bytes), "pdf"

    if ext in (".docx", ".doc"):
        return _docx(file_bytes), "docx"

    if ext in (".pptx", ".ppt"):
        return _pptx(file_bytes), "pptx"

    if ext in TEXT_EXTS:
        return file_bytes.decode("utf-8", errors="replace"), "text"

    if ext in IMAGE_EXTS:
        mime_map = {
            ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp",
        }
        mime = mime_map.get(ext, "image/jpeg")
        b64 = base64.b64encode(file_bytes).decode()
        payload = json.dumps({"type": "image", "mime": mime, "b64": b64})
        return payload, "image"

    if ext in VIDEO_EXTS:
        return (
            f"[Video file '{filename}' received. Automatic video transcription is not supported in this version. "
            f"Please paste a transcript or meeting notes from this recording as a text source instead.]",
            "video",
        )

    # Unknown — try UTF-8
    try:
        return file_bytes.decode("utf-8", errors="replace"), "text"
    except Exception:
        return f"[Binary file '{filename}' — content could not be extracted]", "binary"


def _pdf(file_bytes: bytes) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text.strip())
        return "\n\n".join(parts) if parts else "[PDF: no extractable text found]"
    except ImportError:
        return "[PDF: pdfplumber not installed — run pip install pdfplumber]"
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def _docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # Also grab table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return "\n".join(paragraphs) if paragraphs else "[DOCX: no text found]"
    except ImportError:
        return "[DOCX: python-docx not installed — run pip install python-docx]"
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def _pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
            if slide_lines:
                lines.append(f"[Slide {i}]")
                lines.extend(slide_lines)
        return "\n".join(lines) if lines else "[PPTX: no text found]"
    except ImportError:
        return "[PPTX: python-pptx not installed — run pip install python-pptx]"
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def build_source_meta(filename: str, file_type: str, label: str) -> dict:
    """Build a source metadata dict."""
    return {
        "filename": filename,
        "file_type": file_type,
        "reliability_label": label,
        "reliability_score": RELIABILITY_SCORES.get(label, 6),
    }
